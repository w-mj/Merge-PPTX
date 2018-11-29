import pptx, os, copy, six

def main(origin_list=None):
    if origin_list is None:
        origin_list = os.listdir('.')
        origin_list = [x for x in origin_list if x[-4:] == 'pptx' and x != "merged.pptx"]
        origin_list.sort()
    
    merged = pptx.Presentation()
    for name in origin_list:
        if os.path.isfile(name) and name[-4:] != 'pptx':
            continue
        print('open {}'.format(name))
        input_f = open(name, 'rb')
        in_pre = pptx.Presentation(input_f)
        input_f.close()
        for i in range(len(in_pre.slides)):
            template = in_pre.slides[i]
            try:
                blank_slide_layout = in_pre.slide_layouts[6]
            except:
                blank_slide_layout = in_pre.slide_layouts[len(in_pre.slide_layouts) - 1]

            copied_slide = merged.slides.add_slide(blank_slide_layout)
            for shp in template.shapes:
                if shp.has_text_frame:
                    for par in shp.text_frame.paragraphs:
                        for run in par.runs:
                            font = run.font
                            font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
                el = shp.element
                newel = copy.deepcopy(el)
                copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

            # for _, value in six.iteritems(template.part.rels):
            #     # Make sure we don't copy a notesSlide relation as that won't exist
            #     if "notesSlide" not in value.reltype:
            #         copied_slide.part.rels.add_relationship(value.reltype,
            #                                         value._target,
            #                                         value.rId)

    
    merged.save('merged.pptx')

if __name__ == '__main__':
    main()
        